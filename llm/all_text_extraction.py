import os
import io
import re
from pptx import Presentation
from PIL import Image
import pytesseract

# --- CONFIGURATION ---
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

def clean_text(text: str) -> str:
    """Removes unwanted copyright lines and extra whitespace."""
    pattern = r"© Tata Community Initiatives Trust, all rights reserved\."
    return re.sub(pattern, "", text).strip()

# --- MODIFIED: Function now accepts a file handle instead of a path ---
def extract_ppt_content(pptx_path, out_file, source_filename):
    """Extracts all content from a single PPTX file and writes to an open file handle."""
    try:
        prs = Presentation(pptx_path)
        print(f"✅ Processing: {source_filename}")

        # Writes to the passed 'out_file' handle instead of opening a new file
        for i, slide in enumerate(prs.slides, start=1):
            slide_content = []

            # --- Extract visible text from shapes ---
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = clean_text(shape.text)
                    if text:
                        slide_content.append(text)

            # --- Extract images and run OCR ---
            for shape in slide.shapes:
                if hasattr(shape, 'image'):
                    image = shape.image
                    img_bytes = image.blob
                    try:
                        img = Image.open(io.BytesIO(img_bytes))
                        ocr_text = pytesseract.image_to_string(img)
                        text = clean_text(ocr_text)
                        if text:
                            slide_content.append(f"[[OCR Image Text: {text}]]")
                    except Exception as img_e:
                        print(f"   - Could not process an image on slide {i}: {img_e}")

            # --- Extract speaker notes ---
            notes_text = ""
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                if text_frame and text_frame.text.strip():
                    notes_text = clean_text(text_frame.text)

            # --- Write content with the source filename in the separator ---
            out_file.write(f"\n--- Source: {source_filename}, Slide: {i} ---\n")
            if slide_content:
                out_file.write("\n\n".join(slide_content))
                out_file.write("\n")
            if notes_text:
                out_file.write("\n📝 Notes:\n")
                out_file.write(notes_text)
                out_file.write("\n")

    except Exception as e:
        print(f"❌ An error occurred while processing {source_filename}: {e}")

def process_folder(folder_path, output_filename="ALL14TEXT.txt"):
    """Finds all .pptx files, processes them, and saves all content to a single file."""
    print(f"--- Starting processing for folder: {folder_path} ---\n")

    # --- MODIFIED: Define and open the single output file before the loop ---
    output_filepath = os.path.join(folder_path, output_filename)
    found_pptx = False

    with open(output_filepath, "w", encoding="utf-8") as main_out_file:
        for filename in os.listdir(folder_path):
            if filename.lower().endswith(".pptx"):
                found_pptx = True
                input_pptx_path = os.path.join(folder_path, filename)

                # --- MODIFIED: Pass the open file handle to the extraction function ---
                extract_ppt_content(input_pptx_path, main_out_file, filename)
                print("-" * 20)

    if not found_pptx:
        print("No .pptx files found in the specified folder.")
    else:
        # --- MODIFIED: Final success message refers to the single file ---
        print(f"\n🎉 All presentations processed. Combined output saved to {output_filepath}")


# --- MAIN EXECUTION ---
if __name__ == "__main__":
    SOURCE_FOLDER = r"source_documents"
    # --- MODIFIED: Pass the desired single output filename ---
    process_folder(SOURCE_FOLDER, output_filename="ALL14TEXT.txt")

